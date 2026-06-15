from math import cos, radians, sin
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUT_DIR = Path(r"E:\dotnet-agent-playbook\docs\ai-agent\Agent-Framework\outputs\manual-20260612-rubber-duck-flow\presentations\rubber-duck-flow\output")
PPTX_PATH = OUT_DIR / "rubber-duck-agent-flow-full-loop.pptx"
PNG_PATH = OUT_DIR / "rubber-duck-agent-flow-full-loop-preview.png"

W_IN, H_IN = 16, 9
PX_SCALE = 120
W_PX, H_PX = int(W_IN * PX_SCALE), int(H_IN * PX_SCALE)

BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(184, 192, 204)
BLUE = RGBColor(0, 174, 239)
BLUE_DARK = RGBColor(0, 112, 176)
YELLOW = RGBColor(255, 210, 30)
DARK_PANEL = RGBColor(22, 28, 36)
MID_PANEL = RGBColor(31, 41, 51)
LINE = RGBColor(70, 82, 96)


def set_text(shape, text, size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return shape


def add_text(slide, x, y, w, h, text, size=24, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_text(box, text, size=size, color=color, bold=bold, align=align)
    box.text_frame.margin_left = 0
    box.text_frame.margin_right = 0
    box.text_frame.margin_top = 0
    box.text_frame.margin_bottom = 0
    return box


def add_pill(slide, x, y, w, h, text, fill, line=LINE, text_color=WHITE, size=18, bold=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(1.2)
    set_text(shp, text, size=size, color=text_color, bold=bold, align=PP_ALIGN.CENTER)
    shp.text_frame.margin_left = Inches(0.08)
    shp.text_frame.margin_right = Inches(0.08)
    return shp


def add_card(slide, x, y, w, h, title, subtitle, accent=BLUE):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = DARK_PANEL
    card.line.color.rgb = LINE
    card.line.width = Pt(1)
    add_text(slide, x + 0.26, y + 0.20, w - 0.52, 0.36, title, size=17, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, x + 0.24, y + 0.66, w - 0.48, 0.52, subtitle, size=12.5, color=MUTED, align=PP_ALIGN.CENTER)
    return card


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(W_IN)
    prs.slide_height = Inches(H_IN)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = BLACK

    add_text(slide, 0.85, 0.58, 9.4, 0.56, "/rubber-duck：第二视角审查循环", size=33, bold=True)
    add_text(slide, 0.88, 1.25, 6.2, 0.34, "不是持续聊天的第二助手，而是按需调用的一轮审查器。", size=15, color=MUTED)

    add_pill(slide, 0.88, 1.95, 3.8, 0.52, "单次调用 = 一轮 critique", DARK_PANEL, text_color=WHITE, size=16)
    add_pill(slide, 0.88, 2.62, 3.8, 0.52, "关键节点可反复调用", DARK_PANEL, text_color=YELLOW, size=16)

    add_text(slide, 0.9, 3.58, 4.25, 0.34, "使用时机", size=18, color=WHITE, bold=True)
    add_text(
        slide,
        0.92,
        4.03,
        4.7,
        1.75,
        "写代码前：检查方案是否靠谱\n改代码后：找边界条件和风险\n提交前：补齐测试与真实问题",
        size=16,
        color=MUTED,
    )

    add_text(slide, 0.9, 6.55, 4.8, 0.66, "核心变化：\n从“我说给鸭子听”，变成“另一个模型主动挑刺”。", size=16, color=WHITE)

    # Cycle ring and cards.
    cx, cy = 10.45, 4.73
    # PowerPoint's built-in circular arrow is only a partial arc, so compose a
    # full editable loop from an oval stroke plus a separate arrow head.
    ring = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.53), Inches(2.92), Inches(3.84), Inches(3.84))
    ring.fill.background()
    ring.line.color.rgb = BLUE
    ring.line.width = Pt(15)
    arrow_head = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(11.62), Inches(5.96), Inches(0.58), Inches(0.58))
    arrow_head.rotation = 138
    arrow_head.fill.solid()
    arrow_head.fill.fore_color.rgb = BLUE
    arrow_head.line.color.rgb = BLUE

    add_card(slide, 8.72, 1.55, 3.44, 1.22, "1 生成方案", "Plan / Design / Code / Tests")
    add_card(slide, 12.42, 3.55, 2.82, 1.34, "2 审查问题", "/rubber-duck", accent=YELLOW)
    add_card(slide, 8.92, 6.63, 3.12, 1.24, "3 返回批评", "Blind spots / Risks / Fixes")
    add_card(slide, 6.02, 3.58, 2.72, 1.34, "4 吸收改进", "主 Agent 继续执行")

    # Center concept.
    center = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - 1.05), Inches(cy - 1.05), Inches(2.1), Inches(2.1))
    center.fill.solid()
    center.fill.fore_color.rgb = MID_PANEL
    center.line.color.rgb = BLUE
    center.line.width = Pt(2)
    set_text(center, "Second\nOpinion\n第二视角", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # Tiny code cue.
    code = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.0), Inches(1.03), Inches(2.35), Inches(0.48))
    code.fill.solid()
    code.fill.fore_color.rgb = RGBColor(8, 12, 16)
    code.line.color.rgb = LINE
    set_text(code, "> /rubber-duck", size=14, color=YELLOW, bold=True, align=PP_ALIGN.CENTER, font="Consolas")

    add_text(slide, 6.2, 8.28, 8.9, 0.28, "单次调用返回结构化 critique；主 Agent 吸收反馈后继续计划、实现和测试。", size=13.5, color=MUTED, align=PP_ALIGN.CENTER)

    prs.save(PPTX_PATH)


def try_font(size, bold=False):
    candidates = [
        r"C:\Windows\Fonts\msyhbd.ttc" if bold else r"C:\Windows\Fonts\msyh.ttc",
        r"C:\Windows\Fonts\simhei.ttf",
        r"C:\Windows\Fonts\arial.ttf",
    ]
    for p in candidates:
        if Path(p).exists():
            return ImageFont.truetype(p, size)
    return ImageFont.load_default()


def draw_preview():
    img = Image.new("RGB", (W_PX, H_PX), (0, 0, 0))
    d = ImageDraw.Draw(img)
    S = PX_SCALE

    def xy(x, y):
        return int(x * S), int(y * S)

    def rect(x, y, w, h, fill, outline=(70, 82, 96), radius=16, width=2):
        d.rounded_rectangle([*xy(x, y), *xy(x + w, y + h)], radius=radius, fill=fill, outline=outline, width=width)

    def text(x, y, s, size, fill=(255, 255, 255), bold=False, anchor=None, align="left"):
        font = try_font(size, bold)
        d.multiline_text(xy(x, y), s, font=font, fill=fill, spacing=8, anchor=anchor, align=align)

    text(0.85, 0.58, "/rubber-duck：第二视角审查循环", 54, bold=True)
    text(0.88, 1.25, "不是持续聊天的第二助手，而是按需调用的一轮审查器。", 24, fill=(184, 192, 204))

    rect(0.88, 1.95, 3.8, 0.52, (22, 28, 36))
    text(2.78, 2.08, "单次调用 = 一轮 critique", 22, anchor="ma", bold=True)
    rect(0.88, 2.62, 3.8, 0.52, (22, 28, 36))
    text(2.78, 2.75, "关键节点可反复调用", 22, fill=(255, 210, 30), anchor="ma", bold=True)

    text(0.9, 3.58, "使用时机", 29, bold=True)
    text(0.92, 4.03, "写代码前：检查方案是否靠谱\n改代码后：找边界条件和风险\n提交前：补齐测试与真实问题", 25, fill=(184, 192, 204))
    text(0.9, 6.55, "核心变化：\n从“我说给鸭子听”，变成“另一个模型主动挑刺”。", 25, bold=True)

    cards = [
        (8.72, 1.55, 3.44, 1.22, "1 生成方案", "Plan / Design / Code / Tests", (0, 174, 239)),
        (12.42, 3.55, 2.82, 1.34, "2 审查问题", "/rubber-duck", (255, 210, 30)),
        (8.92, 6.63, 3.12, 1.24, "3 返回批评", "Blind spots / Risks / Fixes", (0, 174, 239)),
        (6.02, 3.58, 2.72, 1.34, "4 吸收改进", "主 Agent 继续执行", (0, 174, 239)),
    ]
    for x, y, w, h, title, sub, accent in cards:
        rect(x, y, w, h, (22, 28, 36))
        text(x + w / 2, y + 0.22, title, 24, fill=accent, bold=True, anchor="ma")
        text(x + w / 2, y + 0.74, sub, 18, fill=(184, 192, 204), anchor="ma", align="center")

    # Draw one aligned full loop arrow to match the editable PowerPoint shapes.
    cx, cy, r = 10.45, 4.73, 1.86
    d.ellipse([*xy(cx - r, cy - r), *xy(cx + r, cy + r)], outline=(0, 174, 239), width=18)
    a = radians(48)
    x, y = cx + r * cos(a), cy + r * sin(a)
    tangent = a + radians(90)
    p1 = xy(x + 0.20 * cos(tangent), y + 0.20 * sin(tangent))
    p2 = xy(x - 0.35 * cos(tangent - 0.50), y - 0.35 * sin(tangent - 0.50))
    p3 = xy(x - 0.35 * cos(tangent + 0.50), y - 0.35 * sin(tangent + 0.50))
    d.polygon([p1, p2, p3], fill=(0, 174, 239))

    d.ellipse([*xy(cx - 1.05, cy - 1.05), *xy(cx + 1.05, cy + 1.05)], fill=(31, 41, 51), outline=(0, 174, 239), width=4)
    text(cx, cy - 0.48, "Second\nOpinion\n第二视角", 28, bold=True, anchor="ma", align="center")
    rect(12.0, 1.03, 2.35, 0.48, (8, 12, 16))
    text(13.175, 1.16, "> /rubber-duck", 22, fill=(255, 210, 30), bold=True, anchor="ma")
    text(10.65, 8.28, "单次调用返回结构化 critique；主 Agent 吸收反馈后继续计划、实现和测试。", 22, fill=(184, 192, 204), anchor="ma")

    img.save(PNG_PATH)


if __name__ == "__main__":
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    build_pptx()
    draw_preview()
    print(PPTX_PATH)
    print(PNG_PATH)
